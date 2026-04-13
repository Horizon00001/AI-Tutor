from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os
import re
import argparse
from pathlib import Path

try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

LATEX_MAP = {
    '\\lceil': '⌈', '\\rceil': '⌉', '\\lfloor': '⌊', '\\rfloor': '⌋',
    '\\log': 'log', '\\infty': '∞', '\\pi': 'π', '\\times': '×',
    '\\div': '÷', '\\leq': '≤', '\\geq': '≥', '\\neq': '≠',
    '\\approx': '≈', '\\pm': '±', '\\cdot': '·', '\\{': '{', '\\}': '}',
    '\\alpha': 'α', '\\beta': 'β', '\\gamma': 'γ', '\\delta': 'δ',
    '\\epsilon': 'ε', '\\theta': 'θ', '\\lambda': 'λ', '\\mu': 'μ',
    '\\sigma': 'σ', '\\omega': 'ω', '\\Delta': 'Δ', '\\Sigma': 'Σ',
    '\\Omega': 'Ω', '\\rightarrow': '→', '\\leftarrow': '←',
    '\\Rightarrow': '⇒', '\\Leftarrow': '⇐', '\\Leftrightarrow': '⇔',
    '\\equiv': '≡', '\\sim': '∼', '\\propto': '∝', '\\subset': '⊂',
    '\\supset': '⊃', '\\in': '∈', '\\ni': '∋', '\\notin': '∉',
    '\\cup': '∪', '\\cap': '∩', '\\emptyset': '∅', '\\forall': '∀',
    '\\exists': '∃', '\\nabla': '∇', '\\partial': '∂', '\\sum': '∑',
    '\\prod': '∏', '\\int': '∫', '\\oint': '∮', '\\sqrt': '√',
    '\\angle': '∠', '\\circ': '∘', '\\bullet': '∙', '\\ast': '∗'
}

def clean_latex(text):
    for k, v in LATEX_MAP.items():
        text = text.replace(k, v)
    return text

def add_math_text(paragraph, text):
    """处理包含 $...$ 或 $$...$$ 的文本，并将其添加到段落中"""
    # 将 $$ 替换为 $，统一处理
    text = text.replace('$$', '$')
    # 按 $ 分割，偶数索引是普通文本，奇数索引是公式
    parts = re.split(r'\$(.*?)\$', text)
    for i, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        if i % 2 == 1: # 公式部分
            run.text = clean_latex(part)
            run.font.name = 'Cambria Math'
            run.font.italic = True
        else: # 普通文本
            run.text = part
            run.font.name = 'Microsoft YaHei'

def parse_options(content):
    """
    尝试从内容中解析 A, B, C, D 选项。
    支持格式：A) A） A. (A) [A]
    """
    options = []
    # 正则表达式匹配常见的选项格式
    # ^\s* 匹配行首空白
    # ([A-D]) 捕获选项字母
    # [)）\.\]] 匹配分隔符（括号、点等）
    # (.*) 捕获选项内容
    option_pattern = re.compile(r'^\s*([A-D])[)）\.\]]\s*(.*)$', re.MULTILINE)
    
    matches = option_pattern.findall(content)
    if matches:
        for letter, text in matches:
            options.append((letter, text.strip()))
    
    return options

def create_title_slide(prs, main_title, subtitle_text):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = main_title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = RGBColor(0, 112, 192) # 主题蓝
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0), Inches(4.0), prs.slide_width, Inches(0.8))
    p2 = subtitle_box.text_frame.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(28)
    p2.font.name = 'Microsoft YaHei'
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER

def create_question_slide(prs, question_num, content, source, answer, analysis):
    slide_layout = prs.slide_layouts[6] # 使用空白版式，自定义文本框
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. 题号标题
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(13.333), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = f"第 {question_num} 题"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = RGBColor(0, 112, 192) # 主题蓝
    p.alignment = PP_ALIGN.CENTER
    
    # 2. 解析题目内容和选项
    options = parse_options(content)
    
    if options:
        # 如果是选择题，第一行通常是题干
        question_text = content.split('\n')[0]
    else:
        # 如果不是选择题，全文都是题干
        question_text = content

    # 3. 题目内容框
    q_height = Inches(1.5) if options else Inches(4.0)
    q_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), q_height)
    q_box.text_frame.word_wrap = True
    p = q_box.text_frame.paragraphs[0]
    add_math_text(p, question_text)
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(14)
    
    # 4. 选项 (如果有)
    if options:
        # 选项框缩进一点
        o_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.333), Inches(1.8))
        o_box.text_frame.word_wrap = True
        
        # 按照每行两个选项进行分组：(A, B) 一行, (C, D) 一行
        for i in range(0, len(options), 2):
            if i == 0:
                p = o_box.text_frame.paragraphs[0]
            else:
                p = o_box.text_frame.add_paragraph()
            
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(12)
            # 添加制表符实现两列对齐 (10.333宽，中间约5.0)
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._p.get_or_add_pPr()
            # 创建 tabStop 元素，val=7620 EMUs ≈ 5.0 英寸
            tab = etree.SubElement(pPr, qn('a:tab'))
            tab.set('val', '7620')
            
            # 处理第一个选项 (A 或 C)
            opt1_letter, opt1_text = options[i]
            add_math_text(p, f"{opt1_letter}) {opt1_text}")
            
            # 如果有第二个选项 (B 或 D)，添加制表符并显示
            if i + 1 < len(options):
                opt2_letter, opt2_text = options[i+1]
                spacing_run = p.add_run()
                spacing_run.text = "\t" 
                add_math_text(p, f"{opt2_letter}) {opt2_text}")
                
            # 设置该行所有文本的字体大小
            for run in p.runs:
                run.font.size = Pt(20)
    
    # 5. 答案与解析 (命名为 AnimatedAnswerShape 以便后续添加动画)
    a_y_pos = Inches(4.6) if options else Inches(4.2)
    a_box = slide.shapes.add_textbox(Inches(1.0), a_y_pos, Inches(11.333), Inches(2.5))
    a_box.name = "AnimatedAnswerShape"
    a_box.text_frame.word_wrap = True
    
    p = a_box.text_frame.paragraphs[0]
    p.text = f"【正确答案】 {answer}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 150, 0) # 绿色
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(10)
    
    p2 = a_box.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    
    run = p2.add_run()
    run.text = "【解析】 "
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(100, 100, 100)
    
    add_math_text(p2, analysis)
    
    for run in p2.runs:
        if not run.font.size:
            run.font.size = Pt(18)

def add_animations_via_com(pptx_path):
    if not HAS_WIN32COM:
        print("未安装 pywin32，跳过添加动画步骤。")
        return
    
    print(f"正在通过 COM 接口为 {pptx_path} 添加动画...")
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        # 将路径转换为绝对路径
        abs_path = os.path.abspath(pptx_path)
        # 以隐藏窗口方式打开
        prs = app.Presentations.Open(abs_path, WithWindow=False)
        
        for slide in prs.Slides:
            for shape in slide.Shapes:
                if shape.Name == "AnimatedAnswerShape":
                    # 10 = msoAnimEffectAppear, 1 = msoAnimTriggerOnPageClick
                    slide.TimeLine.MainSequence.AddEffect(shape, 10, 0, 1)
        
        prs.Save()
        prs.Close()
        app.Quit()
        print("动画添加完成！")
    except Exception as e:
        print(f"添加动画时出现异常: {e}")
        try:
            app.Quit()
        except:
            pass

def main():
    parser = argparse.ArgumentParser(description="根据清理后的 JSON 数据生成 PPT 试卷讲解。")
    parser.add_argument("json_path", help="输入清理后的 JSON 文件路径。")
    parser.add_argument("-o", "--output", help="输出 PPTX 文件路径。如果未提供，将在 output 目录生成。")
    parser.add_argument("-t", "--title", default=None, help="PPT 标题。默认使用JSON文件名。")
    
    args = parser.parse_args()
    
    json_path = Path(args.json_path)
    if not json_path.exists():
        print(f"错误: 找不到输入文件 {json_path}")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    if isinstance(data, list):
        questions = data
        source = None
    else:
        questions = data.get('questions', data)
        source = data.get('source')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    subtitle_text = f"共 {len(questions)} 道题目"
    default_title = json_path.stem.replace('_content_list_cleaned', '').replace('_cleaned', '')
    
    if args.title:
        main_title = args.title
    elif source:
        main_title = source
    elif questions and questions[0].get('source'):
        main_title = questions[0]['source']
    else:
        main_title = default_title
    
    create_title_slide(prs, main_title, subtitle_text)
    
    for idx, q in enumerate(questions, 1):
        create_question_slide(prs, idx, q['content'], q.get('source', ''), q['answer'], q['analysis'])
    
    if args.output:
        output_path = Path(args.output)
    else:
        # 默认输出到 output/ppt 目录，并以输入文件名命名
        output_dir = Path("output/ppt")
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / f"{json_path.stem}_fixed.pptx"
    
    prs.save(str(output_path))
    print(f"PPT已生成：{output_path}")
    
    # 后处理添加动画
    add_animations_via_com(str(output_path))

if __name__ == "__main__":
    main()
