import json
import re
import os
from pathlib import Path
from typing import Literal
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

from services.latex_renderer import latex_renderer

LATEX_MAP = {
    '\\lceil': '⌈', '\\rceil': '⌉', '\\lfloor': '⌊', '\\rfloor': '⌋',
    '\\log': 'log', '\\infty': '∞', '\\pi': 'π', '\\times': '×',
    '\\div': '÷', '\\leq': '≤', '\\geq': '≥', '\\neq': '≠',
    '\\approx': '≈', '\\pm': '±', '\\cdot': '·', '\\{': '{', '\\}': '}',
    '\\alpha': 'α', '\\beta': 'β', '\\gamma': 'γ', '\\delta': 'δ',
    '\\epsilon': 'ε', '\\theta': 'θ', '\\lambda': 'λ', '\\mu': 'μ',
    '\\sigma': 'σ', '\\omega': 'ω', '\\Delta': 'Δ', '\\Sigma': 'Σ',
    '\\Omega': 'Ω', '\\rightarrow': '→', '\\leftarrow': '←',
    '\\Rightarrow': '⇒', '\\Leftarrow': '⇐', '\\Leftrightarrow': '⇔',
    '\\equiv': '≡', '\\sim': '∼', '\\propto': '∝', '\\subset': '⊂',
    '\\supset': '⊃', '\\in': '∈', '\\ni': '∋', '\\notin': '∉',
    '\\cup': '∪', '\\cap': '∩', '\\emptyset': '∅', '\\forall': '∀',
    '\\exists': '∃', '\\nabla': '∇', '\\partial': '∂', '\\sum': '∑',
    '\\prod': '∏', '\\int': '∫', '\\oint': '∮', '\\sqrt': '√',
    '\\angle': '∠', '\\circ': '∘', '\\bullet': '∙', '\\ast': '∗'
}


class PPTGenerator:
    def __init__(self):
        self.latex_map = LATEX_MAP

    def clean_latex(self, text: str) -> str:
        for k, v in self.latex_map.items():
            text = text.replace(k, v)
        return text

    def contains_latex(self, text: str) -> bool:
        """检测文本是否包含 LaTeX 公式（$...$ 包裹）"""
        return '$' in text

    def parse_latex_segments(self, text: str) -> list[tuple[Literal['text', 'latex'], str]]:
        """
        解析文本为片段列表

        Returns:
            [('text', '普通文本'), ('latex', '公式内容'), ...]
        """
        segments = []
        # 转义 \$ 情况先处理
        text = text.replace('\\$', '\x00')

        # 匹配 $...$ 包裹的 LaTeX
        pattern = re.compile(r'\$([^\$]+)\$')
        last_end = 0

        for match in pattern.finditer(text):
            start = match.start()
            if start > last_end:
                plain = text[last_end:start].replace('\x00', '$')
                if plain:
                    segments.append(('text', plain))

            latex_content = match.group(1)
            segments.append(('latex', latex_content))
            last_end = match.end()

        if last_end < len(text):
            plain = text[last_end:].replace('\x00', '$')
            if plain:
                segments.append(('text', plain))

        return segments

    def calculate_content_fit(
        self,
        question_text: str,
        options: list,
        analysis: str,
        content_width: Inches,
        available_height: Inches
    ) -> dict:
        """
        计算内容是否适合幻灯片，返回推荐的字体大小和布局参数
        """
        # 初始字体大小
        title_font_size = 22
        option_font_size = 20
        analysis_font_size = 18

        # 检查是否包含 LaTeX
        has_latex_question = self.contains_latex(question_text)
        has_latex_analysis = self.contains_latex(analysis)

        # 估算各部分内容高度
        def estimate_total_height(q_font, o_font, a_font):
            height = Inches(0.8)  # 标题高度

            # 题目高度 - 考虑 LaTeX
            if has_latex_question:
                # LaTeX 公式通常需要更多空间
                q_height = self.estimate_text_height(question_text, content_width, q_font)
                # 有 LaTeX 时增加额外空间
                q_height = q_height * 1.3
                if options:
                    q_height = min(q_height, Inches(2.0))  # 有选项时限制高度
                else:
                    q_height = min(q_height, Inches(4.0))  # 无选项时可以更高
            else:
                q_height = self.estimate_text_height(question_text, content_width, q_font)
                if options:
                    q_height = min(q_height, Inches(1.5))  # 有选项时限制高度
                else:
                    q_height = min(q_height, Inches(3.5))  # 无选项时可以更高
            height += q_height + Inches(0.2)

            # 选项高度
            if options:
                # 选项数量不同，高度也不同
                num_options = len(options)
                option_height = Inches(0.45 * ((num_options + 1) // 2))  # 每对选项0.45英寸
                height += min(option_height, Inches(2.0))  # 最大2英寸

            # 答案高度
            height += Inches(0.4) + Inches(0.1)

            # 解析高度 - 考虑 LaTeX
            if has_latex_analysis:
                # LaTeX 公式需要更多空间
                a_height = self.estimate_text_height('【解析】 ' + analysis, content_width, a_font)
                a_height = a_height * 1.3  # 增加额外空间
                height += min(a_height, Inches(2.5))  # 最大2.5英寸
            else:
                a_height = self.estimate_text_height('【解析】 ' + analysis, content_width, a_font)
                height += min(a_height, Inches(2.0))  # 最大2英寸

            return height

        # 检查是否需要缩小字体
        total_height = estimate_total_height(title_font_size, option_font_size, analysis_font_size)

        # 如果内容超出，逐步缩小字体
        min_title_font = 12
        min_analysis_font = 10
        
        while total_height > available_height and (
            title_font_size > min_title_font or 
            analysis_font_size > min_analysis_font
        ):
            # 优先缩小解析字体
            if analysis_font_size > min_analysis_font:
                analysis_font_size -= 1
            # 然后缩小标题字体
            elif title_font_size > min_title_font:
                title_font_size -= 1
                option_font_size = max(14, title_font_size - 2)

            total_height = estimate_total_height(title_font_size, option_font_size, analysis_font_size)

        return {
            'title_font_size': title_font_size,
            'option_font_size': option_font_size,
            'analysis_font_size': analysis_font_size,
            'fits': total_height <= available_height,
            'estimated_height': total_height,
            'has_latex_question': has_latex_question,
            'has_latex_analysis': has_latex_analysis
        }

    def estimate_text_height(self, text: str, width: Inches, font_size: int) -> Inches:
        """
        估算文本所需的高度
        基于字符数和字体大小进行估算
        """
        if not text:
            return Inches(0.3)

        # 估算每行可容纳的字符数（基于字体大小和宽度）
        avg_char_width = font_size * 0.015  # 英寸/字符
        chars_per_line = int(width / Inches(avg_char_width))
        chars_per_line = max(chars_per_line, 20)  # 最少每行20个字符

        # 计算需要的行数
        text_length = len(text)
        num_lines = max(1, (text_length + chars_per_line - 1) // chars_per_line)

        # 每行高度（根据字体大小）
        line_height = font_size * 0.025  # 英寸
        total_height = num_lines * line_height + 0.1  # 添加一些padding

        return Inches(total_height)

    def add_mixed_content_to_slide(
        self,
        slide,
        segments: list[tuple[Literal['text', 'latex'], str]],
        x: Inches,
        y: Inches,
        width: Inches,
        height: Inches,
        font_size: int = 20,
        font_bold: bool = False,
    ) -> int:
        """
        向 slide 添加混合内容（文本 + LaTeX 图片）

        Returns:
            使用的 y 位置（用于连续添加内容）
        """
        from PIL import Image as PILImage

        current_y = y
        # 根据字体大小自适应行高
        line_height = Inches(font_size * 0.035)
        # 根据字体大小自适应公式图片高度
        base_font_size = 20
        base_img_height = Inches(0.5)
        max_img_height = base_img_height * (font_size / base_font_size) * 1.2
        max_img_width = width

        for seg_type, seg_content in segments:
            if seg_type == 'text':
                # 清理 LaTeX 字符映射
                clean_text = self.clean_latex(seg_content)

                if not clean_text.strip():
                    continue

                # 估算文本高度
                est_height = self.estimate_text_height(clean_text, width, font_size)

                # 添加文本框
                text_box = slide.shapes.add_textbox(x, current_y, width, est_height)
                tf = text_box.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                p = tf.paragraphs[0]
                p.text = clean_text
                p.font.size = Pt(font_size)
                p.font.bold = font_bold
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.LEFT

                current_y += est_height + Inches(0.05)

            else:  # latex
                # 渲染 LaTeX 为图片
                img_path = latex_renderer.render_to_image(seg_content)

                # 获取图片原始尺寸，保持宽高比
                with PILImage.open(img_path) as img:
                    img_w, img_h = img.size

                # 计算缩放后的尺寸（保持宽高比）
                aspect_ratio = img_w / img_h
                # 根据字体大小计算合适的显示宽度
                display_width = min(max_img_width, Inches(6.0))
                display_height = display_width / aspect_ratio

                # 如果计算的高度超过最大高度，则按高度缩放
                if display_height > max_img_height:
                    display_height = max_img_height
                    display_width = display_height * aspect_ratio

                # 如果宽度超过最大宽度，则按宽度缩放
                if display_width > max_img_width:
                    display_width = max_img_width
                    display_height = display_width / aspect_ratio

                # 添加图片到幻灯片
                pic = slide.shapes.add_picture(img_path, x, current_y, width=display_width)
                current_y += display_height + Inches(0.1)

        return current_y

    def add_math_text(self, paragraph, text: str):
        """保留原有方法用于简单文本（无 $...$）"""
        text = text.replace('$$', '$')
        parts = re.split(r'\$(.*?)\$', text)
        for i, part in enumerate(parts):
            if not part:
                continue
            run = paragraph.add_run()
            if i % 2 == 1:
                # 如果包含复杂 LaTeX，使用清理后的简单版本
                run.text = self.clean_latex(part)
                run.font.name = 'Times New Roman'
                run.font.italic = True
            else:
                run.text = self.clean_latex(part)
                run.font.name = 'Microsoft YaHei'

    def parse_options(self, content: str) -> list:
        option_pattern = re.compile(r'^\s*([A-D])[)）\.\]]\s*(.*)$', re.MULTILINE)
        matches = option_pattern.findall(content)
        options = []
        if matches:
            for letter, text in matches:
                options.append((letter, text.strip()))
        return options

    def create_title_slide(self, prs, main_title: str, subtitle_text: str):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.2))
        p = title_box.text_frame.paragraphs[0]
        p.text = main_title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0), Inches(4.0), prs.slide_width, Inches(0.8))
        p2 = subtitle_box.text_frame.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(28)
        p2.font.name = 'Microsoft YaHei'
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER

    def create_question_slide(self, prs, question_num: int, content: str, source: str, answer: str, analysis: str):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 幻灯片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        content_width = slide_width - 2 * Inches(1.0)  # 内容区域宽度
        max_content_y = slide_height - Inches(0.5)  # 底部边距
        available_height = max_content_y - Inches(1.2)  # 可用高度（从题目开始）

        # 解析选项
        options = self.parse_options(content)

        # 题目内容
        if options:
            question_text = content.split('\n')[0]
        else:
            question_text = content

        # 计算最佳字体大小
        fit_params = self.calculate_content_fit(
            question_text, options, analysis,
            content_width, available_height
        )

        q_font_size = fit_params['title_font_size']
        o_font_size = fit_params['option_font_size']
        a_font_size = fit_params['analysis_font_size']
        has_latex_question = fit_params['has_latex_question']
        has_latex_analysis = fit_params['has_latex_analysis']

        if not fit_params['fits']:
            print(f"警告：第 {question_num} 题内容较长，已自动缩小字体至 {q_font_size}pt")

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0), Inches(0.2), slide_width, Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = f"第 {question_num} 题"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = 'Microsoft YaHei'
        p.font.color.rgb = RGBColor(0, 112, 192)
        p.alignment = PP_ALIGN.CENTER

        current_y = Inches(1.2)

        # 题目内容 - 考虑 LaTeX
        if has_latex_question:
            # 使用混合内容渲染
            segments = self.parse_latex_segments(question_text)
            current_y = self.add_mixed_content_to_slide(
                slide, segments,
                x=Inches(1.0), y=current_y,
                width=content_width, height=Inches(4.0),  # 更大的高度限制
                font_size=q_font_size, font_bold=True
            )
        else:
            # 普通文本 - 估算高度
            q_height = self.estimate_text_height(question_text, content_width, q_font_size)
            if options:
                q_height = min(q_height, Inches(1.5))  # 有选项时限制高度
            else:
                q_height = min(q_height, Inches(3.5))  # 无选项时可以更高

            q_box = slide.shapes.add_textbox(Inches(1.0), current_y, content_width, q_height)
            q_box.text_frame.word_wrap = True
            q_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = q_box.text_frame.paragraphs[0]
            self.add_math_text(p, question_text)
            p.font.size = Pt(q_font_size)
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(14)

            current_y += q_height + Inches(0.2)

        # 选项（如果存在）
        if options:
            option_y = current_y + Inches(0.1)
            # 根据选项数量调整高度
            num_options = len(options)
            option_height = Inches(0.45 * ((num_options + 1) // 2))
            o_box = slide.shapes.add_textbox(Inches(1.5), option_y, content_width - Inches(0.5), option_height)
            o_box.text_frame.word_wrap = True
            o_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            for i in range(0, len(options), 2):
                if i == 0:
                    p = o_box.text_frame.paragraphs[0]
                else:
                    p = o_box.text_frame.add_paragraph()

                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)
                pPr = p._p.get_or_add_pPr()
                tab = etree.SubElement(pPr, qn('a:tab'))
                tab.set('val', '7620')

                opt1_letter, opt1_text = options[i]
                self.add_math_text(p, f"{opt1_letter}) {opt1_text}")

                if i + 1 < len(options):
                    opt2_letter, opt2_text = options[i+1]
                    spacing_run = p.add_run()
                    spacing_run.text = "\t"
                    self.add_math_text(p, f"{opt2_letter}) {opt2_text}")

                for run in p.runs:
                    run.font.size = Pt(o_font_size)

            current_y = option_y + option_height + Inches(0.1)

        # 计算答案和解析区域的位置
        min_answer_y = Inches(4.0)
        a_y_pos = max(current_y + Inches(0.2), min_answer_y)

        # 确保答案区域不会超出边界
        remaining_space = max_content_y - a_y_pos
        if remaining_space < Inches(1.0):
            # 空间严重不足，强制调整位置
            a_y_pos = max_content_y - Inches(1.5)
            print(f"警告：第 {question_num} 题解析区域空间不足，已调整位置")

        # 答案
        answer_height = Inches(0.4)
        answer_box = slide.shapes.add_textbox(Inches(1.0), a_y_pos, content_width, answer_height)
        answer_box.name = "AnimatedAnswerShape"
        answer_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        ap = answer_box.text_frame.paragraphs[0]
        ap.text = f"【正确答案】 {answer}"
        ap.font.size = Pt(20)
        ap.font.bold = True
        ap.font.color.rgb = RGBColor(0, 150, 0)
        ap.alignment = PP_ALIGN.LEFT

        # 解析
        analysis_y = a_y_pos + answer_height + Inches(0.1)
        analysis_segments = self.parse_latex_segments(analysis)

        # 计算解析区域可用高度
        analysis_available_height = max_content_y - analysis_y

        # 确保解析区域有足够空间
        if analysis_available_height < Inches(0.5):
            # 空间严重不足，再次调整
            analysis_y = max_content_y - Inches(1.0)
            analysis_available_height = Inches(0.9)

        if has_latex_analysis:
            # LaTeX 解析需要更多空间
            self.add_mixed_content_to_slide(
                slide, [('text', '【解析】 ')] + analysis_segments,
                x=Inches(1.0), y=analysis_y,
                width=content_width, height=analysis_available_height,
                font_size=a_font_size, font_bold=True
            )
        else:
            # 估算解析文本高度
            analysis_height = self.estimate_text_height('【解析】 ' + analysis, content_width, a_font_size)
            analysis_height = min(analysis_height, analysis_available_height)

            analysis_box = slide.shapes.add_textbox(Inches(1.0), analysis_y, content_width, analysis_height)
            analysis_box.text_frame.word_wrap = True
            analysis_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p2 = analysis_box.text_frame.paragraphs[0]

            run = p2.add_run()
            run.text = "【解析】 "
            run.font.bold = True
            run.font.size = Pt(a_font_size)
            run.font.color.rgb = RGBColor(100, 100, 100)

            self.add_math_text(p2, analysis)

            for run in p2.runs:
                if not run.font.size:
                    run.font.size = Pt(a_font_size)

    def add_animations_via_com(self, pptx_path: Path):
        if not HAS_WIN32COM:
            print("未安装 pywin32，跳过添加动画步骤。")
            return False

        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
            abs_path = os.path.abspath(str(pptx_path))
            prs = app.Presentations.Open(abs_path, WithWindow=False)

            for slide in prs.Slides:
                for shape in slide.Shapes:
                    if shape.Name == "AnimatedAnswerShape":
                        slide.TimeLine.MainSequence.AddEffect(shape, 10, 0, 1)

            prs.Save()
            prs.Close()
            app.Quit()
            return True
        except Exception as e:
            print(f"添加动画时出现异常: {e}")
            try:
                app.Quit()
            except:
                pass
            return False

    def generate(self, questions: list, output_path: Path, title: str = None, use_animation: bool = True) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        subtitle_text = f"共 {len(questions)} 道题目"

        if title:
            main_title = title
        elif questions and questions[0].get('source'):
            main_title = questions[0]['source']
        else:
            main_title = "试卷评讲"

        self.create_title_slide(prs, main_title, subtitle_text)

        for idx, q in enumerate(questions, 1):
            self.create_question_slide(
                prs, idx,
                q.get('content', ''),
                q.get('source', ''),
                q.get('answer', ''),
                q.get('analysis', '')
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        if use_animation:
            self.add_animations_via_com(output_path)

        return output_path


ppt_generator = PPTGenerator()
